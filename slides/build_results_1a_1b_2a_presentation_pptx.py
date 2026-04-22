from __future__ import annotations

import logging
import shutil
import subprocess
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


LOGGER = logging.getLogger(__name__)
SLIDE_WIDTH_INCHES = 13.333333
SLIDE_HEIGHT_INCHES = 7.5


def run_command(command_arguments: list[str], working_directory: Path | None = None) -> subprocess.CompletedProcess[str]:
    """Run a subprocess command and raise a readable error when it fails."""
    completed_process = subprocess.run(
        command_arguments,
        cwd=working_directory,
        check=False,
        capture_output=True,
        text=True,
    )
    if completed_process.returncode != 0:
        raise RuntimeError(
            "Command failed with non-zero exit status.\n"
            f"Command: {' '.join(command_arguments)}\n"
            f"Stdout:\n{completed_process.stdout}\n"
            f"Stderr:\n{completed_process.stderr}"
        )
    return completed_process


def get_pdf_page_count(input_pdf_path: Path) -> int:
    """Read the total page count from pdfinfo so every PDF page becomes one PPT slide."""
    completed_process = run_command(["pdfinfo", str(input_pdf_path)])
    for output_line in completed_process.stdout.splitlines():
        if output_line.startswith("Pages:"):
            return int(output_line.split(":", maxsplit=1)[1].strip())
    raise RuntimeError(f"Could not determine page count for {input_pdf_path}")


def convert_pdf_page_to_wmf(input_pdf_path: Path, page_number: int, working_directory: Path) -> Path:
    """Convert a single PDF page to a WMF file to preserve vector quality inside PowerPoint."""
    svg_like_page_path = working_directory / f"page_{page_number:03d}"
    wmf_output_path = working_directory / f"page_{page_number:03d}.wmf"

    run_command(
        [
            "pdftocairo",
            "-f",
            str(page_number),
            "-l",
            str(page_number),
            "-svg",
            str(input_pdf_path),
            str(svg_like_page_path),
        ]
    )
    run_command(
        [
            "soffice",
            "--headless",
            "--convert-to",
            "wmf",
            "--outdir",
            str(working_directory),
            str(svg_like_page_path),
        ]
    )

    if not wmf_output_path.exists():
        raise RuntimeError(f"Expected WMF output was not created: {wmf_output_path}")
    return wmf_output_path


def build_presentation_from_pdf_pages(input_pdf_path: Path, output_pptx_path: Path) -> None:
    """Create a PowerPoint where each slide is a full-slide EMF rendering of one PDF page."""
    page_count = get_pdf_page_count(input_pdf_path)
    LOGGER.info("Detected %s pages in %s", page_count, input_pdf_path)

    temporary_directory = Path(tempfile.mkdtemp(prefix=f"{output_pptx_path.stem}_build_"))

    try:
        presentation = Presentation()
        presentation.slide_width = Inches(SLIDE_WIDTH_INCHES)
        presentation.slide_height = Inches(SLIDE_HEIGHT_INCHES)
        blank_slide_layout = presentation.slide_layouts[6]

        for page_number in range(1, page_count + 1):
            LOGGER.info("Converting page %s/%s to WMF", page_number, page_count)
            wmf_page_path = convert_pdf_page_to_wmf(input_pdf_path, page_number, temporary_directory)

            slide = presentation.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(
                str(wmf_page_path),
                left=0,
                top=0,
                width=presentation.slide_width,
                height=presentation.slide_height,
            )
        presentation.save(output_pptx_path)
        LOGGER.info("Saved PowerPoint to %s", output_pptx_path)
    finally:
        shutil.rmtree(temporary_directory, ignore_errors=True)


def main() -> None:
    """Build a PPTX version of the compiled Beamer deck while keeping vector page fidelity."""
    logging.basicConfig(
        level=logging.INFO,
        format="%(asctime)s | %(levelname)s | %(message)s",
    )
    slides_directory = Path(__file__).resolve().parent
    input_pdf_path = slides_directory / "results_1a_1b_2a_presentation.pdf"
    output_pptx_path = slides_directory / "results_1a_1b_2a_presentation.pptx"
    build_presentation_from_pdf_pages(input_pdf_path, output_pptx_path)


if __name__ == "__main__":
    main()
