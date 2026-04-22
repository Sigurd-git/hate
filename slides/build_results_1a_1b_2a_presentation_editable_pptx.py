from __future__ import annotations

import json
import logging
import math
import re
import subprocess
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt


LOGGER = logging.getLogger(__name__)

SLIDE_WIDTH_INCHES = 13.333333
SLIDE_HEIGHT_INCHES = 7.5
LEFT_MARGIN_INCHES = 0.55
RIGHT_MARGIN_INCHES = 0.55
TOP_MARGIN_INCHES = 0.22
BOTTOM_MARGIN_INCHES = 0.25
TITLE_HEIGHT_INCHES = 0.62
DIVIDER_Y_INCHES = 0.96
CONTENT_TOP_INCHES = 1.08
CONTENT_WIDTH_INCHES = SLIDE_WIDTH_INCHES - LEFT_MARGIN_INCHES - RIGHT_MARGIN_INCHES
CONTENT_BOTTOM_INCHES = SLIDE_HEIGHT_INCHES - BOTTOM_MARGIN_INCHES
CONTENT_HEIGHT_INCHES = CONTENT_BOTTOM_INCHES - CONTENT_TOP_INCHES

TITLE_FONT_NAME = "Noto Sans CJK SC"
BODY_FONT_NAME = "Noto Sans CJK SC"
MONO_FONT_NAME = "Noto Sans Mono CJK SC"

TITLE_FONT_SIZE_POINTS = 21
BODY_FONT_SIZE_POINTS = 14
SMALL_BODY_FONT_SIZE_POINTS = 12
TABLE_FONT_SIZE_POINTS = 9
CAPTION_FONT_SIZE_POINTS = 11
PAGE_NUMBER_FONT_SIZE_POINTS = 9

TITLE_COLOR = RGBColor(31, 78, 121)
BODY_COLOR = RGBColor(34, 34, 34)
MUTED_COLOR = RGBColor(102, 102, 102)
RULE_COLOR = RGBColor(200, 211, 223)
BLOCK_FILL_COLOR = RGBColor(236, 243, 250)
BLOCK_LINE_COLOR = RGBColor(109, 158, 197)
ALERT_FILL_COLOR = RGBColor(255, 245, 232)
ALERT_LINE_COLOR = RGBColor(229, 160, 92)
TITLE_SLIDE_FILL_COLOR = RGBColor(31, 78, 121)
TITLE_SLIDE_ACCENT_COLOR = RGBColor(129, 168, 202)
WHITE_COLOR = RGBColor(255, 255, 255)


@dataclass
class ImagePanel:
    source_path: Path
    label_text: str = ""


def run_command(command_arguments: list[str], working_directory: Path | None = None) -> subprocess.CompletedProcess[str]:
    """Run a subprocess command and surface stdout/stderr when it fails."""
    completed_process = subprocess.run(
        command_arguments,
        cwd=working_directory,
        check=False,
        capture_output=True,
        text=True,
    )
    if completed_process.returncode != 0:
        raise RuntimeError(
            "Command failed.\n"
            f"Command: {' '.join(command_arguments)}\n"
            f"Stdout:\n{completed_process.stdout}\n"
            f"Stderr:\n{completed_process.stderr}"
        )
    return completed_process


def rgb_fill(shape: Any, color: RGBColor) -> None:
    """Apply a solid fill color to a PowerPoint shape."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def rgb_line(shape: Any, color: RGBColor, width_points: float = 1.0) -> None:
    """Apply a solid line color to a PowerPoint shape."""
    shape.line.color.rgb = color
    shape.line.width = Pt(width_points)


def estimate_wrapped_line_count(text_content: str, characters_per_line: int) -> int:
    """Estimate how many wrapped lines a paragraph will occupy inside a text box."""
    if not text_content.strip():
        return 1
    line_count = 0
    for raw_line in text_content.splitlines():
        cleaned_line = raw_line.strip()
        if not cleaned_line:
            line_count += 1
            continue
        line_count += max(1, math.ceil(len(cleaned_line) / max(characters_per_line, 1)))
    return max(line_count, 1)


def simplify_latex_math(latex_text: str) -> str:
    """Convert the limited math syntax in this deck into readable plain Unicode text."""
    cleaned_text = latex_text
    cleaned_text = cleaned_text.replace("\\!", "")
    cleaned_text = cleaned_text.replace("\\times", "×")
    cleaned_text = cleaned_text.replace("\\to", "→")
    cleaned_text = cleaned_text.replace("\\Rightarrow", "⇒")
    cleaned_text = cleaned_text.replace("\\leftrightarrow", "↔")
    cleaned_text = cleaned_text.replace("\\cdot", "·")
    cleaned_text = cleaned_text.replace("\\pm", "±")
    cleaned_text = cleaned_text.replace("\\rho_s", "ρs")
    cleaned_text = cleaned_text.replace("\\rho", "ρ")
    cleaned_text = cleaned_text.replace("\\Delta", "Δ")
    cleaned_text = cleaned_text.replace("\\approx", "≈")
    cleaned_text = cleaned_text.replace("\\gg", "≫")
    cleaned_text = cleaned_text.replace("\\ll", "≪")
    cleaned_text = cleaned_text.replace("\\geq", "≥")
    cleaned_text = cleaned_text.replace("\\leq", "≤")
    cleaned_text = cleaned_text.replace("\\text", "")
    cleaned_text = cleaned_text.replace("{", "")
    cleaned_text = cleaned_text.replace("}", "")
    cleaned_text = cleaned_text.replace("_", "")
    cleaned_text = cleaned_text.replace("^", "")
    cleaned_text = cleaned_text.replace("\\", "")
    cleaned_text = re.sub(r"\s+", " ", cleaned_text)
    return cleaned_text.strip()


def inlines_to_text(inline_elements: list[dict[str, Any]]) -> str:
    """Flatten a Pandoc inline list into editable plain text."""
    text_fragments: list[str] = []
    for inline_element in inline_elements:
        inline_type = inline_element["t"]
        inline_content = inline_element.get("c")

        if inline_type == "Str":
            text_fragments.append(str(inline_content))
        elif inline_type == "Space":
            text_fragments.append(" ")
        elif inline_type == "SoftBreak":
            text_fragments.append(" ")
        elif inline_type == "LineBreak":
            text_fragments.append("\n")
        elif inline_type in {"Emph", "Strong", "Underline", "Strikeout", "SmallCaps", "Subscript", "Superscript"}:
            text_fragments.append(inlines_to_text(inline_content))
        elif inline_type == "Span":
            text_fragments.append(inlines_to_text(inline_content[1]))
        elif inline_type == "Code":
            text_fragments.append(str(inline_content[1]))
        elif inline_type == "Math":
            text_fragments.append(simplify_latex_math(str(inline_content[1])))
        elif inline_type == "Quoted":
            text_fragments.append(f"\"{inlines_to_text(inline_content[1])}\"")
        elif inline_type == "Note":
            continue
        elif inline_type == "Image":
            continue
        else:
            text_fragments.append(str(inline_content) if inline_content is not None else "")

    flattened_text = "".join(text_fragments)
    flattened_text = re.sub(r"[ \t]+\n", "\n", flattened_text)
    flattened_text = re.sub(r"\n{3,}", "\n\n", flattened_text)
    flattened_text = re.sub(r" {2,}", " ", flattened_text)
    return flattened_text.strip()


def blocks_to_text(block_elements: list[dict[str, Any]]) -> str:
    """Flatten a Pandoc block list into simple multiline text."""
    block_text_parts: list[str] = []
    for block_element in block_elements:
        block_type = block_element["t"]
        if block_type in {"Para", "Plain"}:
            text_value = inlines_to_text(block_element["c"])
            if text_value:
                block_text_parts.append(text_value)
        elif block_type == "BulletList":
            for bullet_level, bullet_text, is_numbered in flatten_list_block(block_element):
                bullet_prefix = f"{bullet_level * '    '}{'1.' if is_numbered else '•'} "
                block_text_parts.append(f"{bullet_prefix}{bullet_text}")
        elif block_type == "OrderedList":
            for bullet_level, bullet_text, is_numbered in flatten_list_block(block_element):
                bullet_prefix = f"{bullet_level * '    '}{'1.' if is_numbered else '•'} "
                block_text_parts.append(f"{bullet_prefix}{bullet_text}")
        elif block_type == "Div":
            block_text_parts.append(blocks_to_text(block_element["c"][1]))
    return "\n".join(text_value for text_value in block_text_parts if text_value.strip())


def flatten_list_block(list_block: dict[str, Any], level: int = 0) -> list[tuple[int, str, bool]]:
    """Flatten nested BulletList or OrderedList nodes into display-ready lines."""
    flattened_items: list[tuple[int, str, bool]] = []
    list_type = list_block["t"]
    is_numbered = list_type == "OrderedList"
    list_items = list_block["c"][1] if is_numbered else list_block["c"]

    for list_item in list_items:
        item_blocks = list_item if not is_numbered else list_item
        paragraph_text_parts: list[str] = []
        for item_block in item_blocks:
            if item_block["t"] in {"Para", "Plain"}:
                paragraph_text_parts.append(inlines_to_text(item_block["c"]))
            elif item_block["t"] in {"BulletList", "OrderedList"}:
                if paragraph_text_parts:
                    flattened_items.append((level, " ".join(paragraph_text_parts).strip(), is_numbered))
                    paragraph_text_parts.clear()
                flattened_items.extend(flatten_list_block(item_block, level + 1))
            elif item_block["t"] == "Div":
                paragraph_text_parts.append(blocks_to_text(item_block["c"][1]))
            else:
                paragraph_text_parts.append(blocks_to_text([item_block]))
        if paragraph_text_parts:
            flattened_items.append((level, " ".join(paragraph_text_parts).strip(), is_numbered))

    return flattened_items


def resolve_image_path(slides_directory: Path, image_reference: str) -> Path:
    """Resolve an image path recorded by Pandoc into an absolute repository path."""
    return (slides_directory / image_reference).resolve()


def extract_image_panels_from_inlines(slides_directory: Path, inline_elements: list[dict[str, Any]]) -> list[ImagePanel]:
    """Extract image panels and optional labels from a Pandoc inline sequence."""
    extracted_panels: list[ImagePanel] = []
    active_image_path: Path | None = None
    active_label_parts: list[str] = []

    def flush_active_panel() -> None:
        nonlocal active_image_path, active_label_parts
        if active_image_path is None:
            return
        extracted_panels.append(ImagePanel(source_path=active_image_path, label_text="".join(active_label_parts).strip()))
        active_image_path = None
        active_label_parts = []

    for inline_element in inline_elements:
        inline_type = inline_element["t"]
        inline_content = inline_element.get("c")

        if inline_type == "Image":
            flush_active_panel()
            active_image_path = resolve_image_path(slides_directory, inline_content[2][0])
        elif inline_type == "SoftBreak":
            flush_active_panel()
        elif inline_type == "LineBreak":
            if active_image_path is not None:
                active_label_parts.append("\n")
        else:
            if active_image_path is not None:
                active_label_parts.append(inlines_to_text([inline_element]))

    flush_active_panel()
    return extracted_panels


def is_image_only_paragraph(block_element: dict[str, Any]) -> bool:
    """Check whether a paragraph contains only a single image object."""
    if block_element["t"] != "Para":
        return False
    non_whitespace_inlines = [inline_element for inline_element in block_element["c"] if inline_element["t"] not in {"Space", "SoftBreak", "LineBreak"}]
    return len(non_whitespace_inlines) == 1 and non_whitespace_inlines[0]["t"] == "Image"


def extract_table_rows(table_block: dict[str, Any]) -> tuple[list[list[str]], str]:
    """Expand a Pandoc table into a simple 2D string grid plus an optional caption."""
    caption_inline_blocks = table_block["c"][1][1]
    table_caption = blocks_to_text(caption_inline_blocks) if caption_inline_blocks else ""

    header_rows = table_block["c"][3][1]
    body_groups = table_block["c"][4]

    expanded_rows: list[list[str]] = []

    def append_row(row_definition: list[Any]) -> None:
        expanded_cells: list[str] = []
        for cell_definition in row_definition[1]:
            column_span = int(cell_definition[3])
            cell_text = blocks_to_text(cell_definition[4])
            expanded_cells.append(cell_text)
            for _ in range(column_span - 1):
                expanded_cells.append("")
        expanded_rows.append(expanded_cells)

    for header_row in header_rows:
        append_row(header_row)

    for body_group in body_groups:
        for body_row in body_group[3]:
            append_row(body_row)

    maximum_column_count = max((len(row_values) for row_values in expanded_rows), default=0)
    for row_values in expanded_rows:
        row_values.extend([""] * (maximum_column_count - len(row_values)))

    return expanded_rows, table_caption


class EditablePresentationBuilder:
    """Build an editable PowerPoint deck from the Beamer source via Pandoc JSON."""

    def __init__(self, slides_directory: Path) -> None:
        self.slides_directory = slides_directory
        self.output_presentation_path = slides_directory / "results_1a_1b_2a_presentation_editable.pptx"
        self.temporary_directory = Path(tempfile.mkdtemp(prefix="results_1a_1b_2a_editable_pptx_"))
        self.image_cache: dict[Path, Path] = {}

    def load_pandoc_document(self) -> dict[str, Any]:
        """Parse the LaTeX source into a Pandoc JSON document."""
        completed_process = run_command(
            [
                "pandoc",
                "-f",
                "latex",
                "-t",
                "json",
                str(self.slides_directory / "results_1a_1b_2a_presentation.tex"),
            ]
        )
        return json.loads(completed_process.stdout)

    def convert_source_figure(self, source_path: Path) -> Path:
        """Convert PDF figures to high-resolution PNGs while reusing native bitmap assets unchanged."""
        resolved_source_path = source_path.resolve()
        if resolved_source_path in self.image_cache:
            return self.image_cache[resolved_source_path]

        suffix = resolved_source_path.suffix.lower()
        if suffix in {".png", ".jpg", ".jpeg"}:
            self.image_cache[resolved_source_path] = resolved_source_path
            return resolved_source_path

        if suffix == ".pdf":
            output_prefix = self.temporary_directory / resolved_source_path.stem
            run_command(
                [
                    "pdftocairo",
                    "-singlefile",
                    "-png",
                    "-r",
                    "400",
                    str(resolved_source_path),
                    str(output_prefix),
                ]
            )
            converted_path = output_prefix.with_suffix(".png")
            self.image_cache[resolved_source_path] = converted_path
            return converted_path

        raise RuntimeError(f"Unsupported image type for {resolved_source_path}")

    def add_contained_picture(self, slide: Any, source_path: Path, x_inches: float, y_inches: float, width_inches: float, height_inches: float) -> None:
        """Insert a picture inside a bounding box while preserving the source aspect ratio."""
        rendered_image_path = self.convert_source_figure(source_path)
        with Image.open(rendered_image_path) as opened_image:
            image_width_pixels, image_height_pixels = opened_image.size

        box_aspect_ratio = width_inches / max(height_inches, 0.01)
        image_aspect_ratio = image_width_pixels / max(image_height_pixels, 1)

        if image_aspect_ratio >= box_aspect_ratio:
            final_width_inches = width_inches
            final_height_inches = width_inches / image_aspect_ratio
            final_x_inches = x_inches
            final_y_inches = y_inches + (height_inches - final_height_inches) / 2
        else:
            final_height_inches = height_inches
            final_width_inches = height_inches * image_aspect_ratio
            final_x_inches = x_inches + (width_inches - final_width_inches) / 2
            final_y_inches = y_inches

        slide.shapes.add_picture(
            str(rendered_image_path),
            Inches(final_x_inches),
            Inches(final_y_inches),
            width=Inches(final_width_inches),
            height=Inches(final_height_inches),
        )

    def add_textbox(self, slide: Any, text_content: str, x_inches: float, y_inches: float, width_inches: float, height_inches: float, *, font_size_points: int, bold: bool = False, italic: bool = False, font_color: RGBColor = BODY_COLOR, alignment: PP_PARAGRAPH_ALIGNMENT = PP_ALIGN.LEFT, fill_color: RGBColor | None = None, line_color: RGBColor | None = None) -> Any:
        """Create a textbox with consistent typography and optional background styling."""
        textbox_shape = slide.shapes.add_textbox(Inches(x_inches), Inches(y_inches), Inches(width_inches), Inches(height_inches))
        text_frame = textbox_shape.text_frame
        text_frame.clear()
        text_frame.margin_left = Pt(4)
        text_frame.margin_right = Pt(4)
        text_frame.margin_top = Pt(3)
        text_frame.margin_bottom = Pt(3)
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        text_frame.word_wrap = True

        if fill_color is not None:
            rgb_fill(textbox_shape, fill_color)
        else:
            textbox_shape.fill.background()

        if line_color is not None:
            rgb_line(textbox_shape, line_color)
        else:
            textbox_shape.line.fill.background()

        paragraphs = text_content.splitlines() if text_content else [""]
        for paragraph_index, paragraph_text in enumerate(paragraphs):
            paragraph = text_frame.paragraphs[0] if paragraph_index == 0 else text_frame.add_paragraph()
            paragraph.text = paragraph_text
            paragraph.alignment = alignment
            paragraph.space_after = Pt(1)
            paragraph.space_before = Pt(0)
            run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
            run.font.name = BODY_FONT_NAME
            run.font.size = Pt(font_size_points)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = font_color

        return textbox_shape

    def render_title_slide(self, presentation: Presentation, metadata: dict[str, Any]) -> None:
        """Render the first slide from the title metadata rather than the empty titlepage frame."""
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = TITLE_SLIDE_FILL_COLOR

        title_text = inlines_to_text(metadata["title"]["c"])
        author_text = "\n".join(inlines_to_text(author_item["c"]) for author_item in metadata["author"]["c"])
        date_text = inlines_to_text(metadata["date"]["c"])

        self.add_textbox(
            slide,
            title_text,
            0.9,
            1.35,
            11.2,
            1.7,
            font_size_points=29,
            bold=True,
            font_color=WHITE_COLOR,
        )

        accent_shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0.9),
            Inches(3.62),
            Inches(2.8),
            Inches(0.05),
        )
        rgb_fill(accent_shape, TITLE_SLIDE_ACCENT_COLOR)
        accent_shape.line.fill.background()

        self.add_textbox(
            slide,
            author_text,
            0.9,
            3.82,
            3.8,
            0.45,
            font_size_points=16,
            font_color=WHITE_COLOR,
        )
        self.add_textbox(
            slide,
            date_text,
            0.9,
            4.32,
            3.8,
            0.38,
            font_size_points=13,
            font_color=TITLE_SLIDE_ACCENT_COLOR,
        )

    def normalize_frame(self, frame: dict[str, Any]) -> tuple[str, list[dict[str, Any]]]:
        """Extract the frame title and normalize the remaining Pandoc blocks."""
        raw_blocks = frame["c"][1]
        frame_title = ""

        if raw_blocks and raw_blocks[0]["t"] == "Para":
            first_inline_elements = raw_blocks[0]["c"]
            if first_inline_elements and first_inline_elements[0]["t"] == "Span":
                frame_title = inlines_to_text(first_inline_elements[0]["c"][1])
                trailing_inline_elements = first_inline_elements[1:]
                while trailing_inline_elements and trailing_inline_elements[0]["t"] in {"Space", "SoftBreak", "LineBreak"}:
                    trailing_inline_elements = trailing_inline_elements[1:]
                normalized_blocks = ([{"t": "Para", "c": trailing_inline_elements}] if trailing_inline_elements else []) + raw_blocks[1:]
            else:
                normalized_blocks = raw_blocks
        else:
            normalized_blocks = raw_blocks

        parsed_blocks: list[dict[str, Any]] = []
        block_index = 0
        while block_index < len(normalized_blocks):
            block_element = normalized_blocks[block_index]

            if is_image_only_paragraph(block_element):
                grouped_image_paths: list[ImagePanel] = []
                while block_index < len(normalized_blocks) and is_image_only_paragraph(normalized_blocks[block_index]):
                    grouped_image_paths.extend(extract_image_panels_from_inlines(self.slides_directory, normalized_blocks[block_index]["c"]))
                    block_index += 1
                parsed_blocks.append({"kind": "image_grid", "rows": self.chunk_image_panels(grouped_image_paths)})
                continue

            parsed_blocks.append(self.parse_single_block(block_element))
            block_index += 1

        cleaned_blocks = [block_entry for block_entry in parsed_blocks if block_entry.get("kind") != "empty"]
        merged_blocks: list[dict[str, Any]] = []
        for block_entry in cleaned_blocks:
            if merged_blocks and block_entry["kind"] == "image_grid" and merged_blocks[-1]["kind"] == "image_grid":
                merged_blocks[-1]["rows"].extend(block_entry["rows"])
            else:
                merged_blocks.append(block_entry)

        return frame_title, merged_blocks

    def chunk_image_panels(self, image_panels: list[ImagePanel]) -> list[list[ImagePanel]]:
        """Group a flat image list into a sensible grid."""
        if not image_panels:
            return []

        panel_count = len(image_panels)
        if panel_count in {2, 4}:
            columns_per_row = 2
        elif panel_count in {3, 6, 9}:
            columns_per_row = 3
        else:
            columns_per_row = min(3, panel_count)

        return [image_panels[index:index + columns_per_row] for index in range(0, panel_count, columns_per_row)]

    def parse_callout_div(self, div_block: dict[str, Any], alert: bool) -> dict[str, Any]:
        """Convert block/alertblock content into a callout descriptor."""
        child_blocks = div_block["c"][1]
        title_text = ""
        body_parts: list[str] = []

        if child_blocks and child_blocks[0]["t"] == "Para":
            first_paragraph_inlines = child_blocks[0]["c"]
            if first_paragraph_inlines and first_paragraph_inlines[0]["t"] == "Span":
                title_text = inlines_to_text(first_paragraph_inlines[0]["c"][1])
                body_parts.append(inlines_to_text(first_paragraph_inlines[1:]))
                child_blocks = child_blocks[1:]

        for child_block in child_blocks:
            body_parts.append(blocks_to_text([child_block]))

        return {
            "kind": "callout",
            "title": title_text.strip(),
            "body": "\n".join(part for part in body_parts if part.strip()).strip(),
            "alert": alert,
        }

    def parse_single_block(self, block_element: dict[str, Any]) -> dict[str, Any]:
        """Convert one Pandoc block into a simpler rendering descriptor."""
        block_type = block_element["t"]

        if block_type == "Para":
            extracted_panels = extract_image_panels_from_inlines(self.slides_directory, block_element["c"])
            if extracted_panels:
                return {"kind": "image_grid", "rows": [extracted_panels]}

            paragraph_text = inlines_to_text(block_element["c"])
            return {"kind": "paragraph", "text": paragraph_text} if paragraph_text else {"kind": "empty"}

        if block_type in {"BulletList", "OrderedList"}:
            return {"kind": "list", "items": flatten_list_block(block_element)}

        if block_type == "Table":
            table_rows, table_caption = extract_table_rows(block_element)
            return {"kind": "table", "rows": table_rows, "caption": table_caption}

        if block_type == "Div":
            div_classes = set(block_element["c"][0][1])
            if "columns" in div_classes:
                column_children = block_element["c"][1]
                if any(child["t"] == "Div" and ("block" in child["c"][0][1] or "alertblock" in child["c"][0][1]) for child in column_children):
                    parsed_columns = []
                    for child in column_children:
                        child_classes = set(child["c"][0][1]) if child["t"] == "Div" else set()
                        if "alertblock" in child_classes:
                            parsed_columns.append(self.parse_callout_div(child, alert=True))
                        elif "block" in child_classes:
                            parsed_columns.append(self.parse_callout_div(child, alert=False))
                    return {"kind": "callout_columns", "columns": parsed_columns}

                grid_rows: list[list[ImagePanel]] = []
                for child in column_children:
                    if child["t"] == "Para":
                        child_panels = extract_image_panels_from_inlines(self.slides_directory, child["c"])
                        if child_panels:
                            grid_rows.append(child_panels)
                return {"kind": "image_grid", "rows": grid_rows} if grid_rows else {"kind": "empty"}

            if "alertblock" in div_classes:
                return self.parse_callout_div(block_element, alert=True)
            if "block" in div_classes:
                return self.parse_callout_div(block_element, alert=False)

            nested_text = blocks_to_text(block_element["c"][1])
            return {"kind": "paragraph", "text": nested_text} if nested_text else {"kind": "empty"}

        nested_text = blocks_to_text([block_element])
        return {"kind": "paragraph", "text": nested_text} if nested_text else {"kind": "empty"}

    def add_standard_slide_chrome(self, slide: Any, frame_title: str, frame_number: int, total_frames: int) -> None:
        """Add the consistent title, divider, and page number elements to a content slide."""
        self.add_textbox(
            slide,
            frame_title,
            LEFT_MARGIN_INCHES,
            TOP_MARGIN_INCHES,
            CONTENT_WIDTH_INCHES,
            TITLE_HEIGHT_INCHES,
            font_size_points=TITLE_FONT_SIZE_POINTS,
            bold=True,
            font_color=TITLE_COLOR,
        )

        divider_shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(LEFT_MARGIN_INCHES),
            Inches(DIVIDER_Y_INCHES),
            Inches(CONTENT_WIDTH_INCHES),
            Inches(0.018),
        )
        rgb_fill(divider_shape, RULE_COLOR)
        divider_shape.line.fill.background()

        page_number_shape = self.add_textbox(
            slide,
            f"{frame_number}/{total_frames}",
            SLIDE_WIDTH_INCHES - 1.1,
            SLIDE_HEIGHT_INCHES - 0.27,
            0.65,
            0.18,
            font_size_points=PAGE_NUMBER_FONT_SIZE_POINTS,
            font_color=MUTED_COLOR,
            alignment=PP_ALIGN.RIGHT,
        )
        page_number_shape.fill.background()

    def render_callout(self, slide: Any, callout_descriptor: dict[str, Any], x_inches: float, y_inches: float, width_inches: float, height_inches: float) -> None:
        """Render a block or alertblock as an editable rounded rectangle."""
        fill_color = ALERT_FILL_COLOR if callout_descriptor["alert"] else BLOCK_FILL_COLOR
        line_color = ALERT_LINE_COLOR if callout_descriptor["alert"] else BLOCK_LINE_COLOR
        body_text = callout_descriptor["body"]
        combined_text = f"{callout_descriptor['title']}\n{body_text}".strip() if callout_descriptor["title"] else body_text

        box_shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x_inches),
            Inches(y_inches),
            Inches(width_inches),
            Inches(height_inches),
        )
        rgb_fill(box_shape, fill_color)
        rgb_line(box_shape, line_color, 1.1)

        text_frame = box_shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.margin_left = Pt(6)
        text_frame.margin_right = Pt(6)
        text_frame.margin_top = Pt(4)
        text_frame.margin_bottom = Pt(4)

        if callout_descriptor["title"]:
            title_paragraph = text_frame.paragraphs[0]
            title_paragraph.text = callout_descriptor["title"]
            title_paragraph.alignment = PP_ALIGN.LEFT
            title_run = title_paragraph.runs[0]
            title_run.font.name = BODY_FONT_NAME
            title_run.font.size = Pt(BODY_FONT_SIZE_POINTS)
            title_run.font.bold = True
            title_run.font.color.rgb = TITLE_COLOR

            if body_text:
                body_paragraph = text_frame.add_paragraph()
                body_paragraph.text = body_text
                body_paragraph.alignment = PP_ALIGN.LEFT
                body_paragraph.space_before = Pt(2)
                body_run = body_paragraph.runs[0]
                body_run.font.name = BODY_FONT_NAME
                body_run.font.size = Pt(SMALL_BODY_FONT_SIZE_POINTS)
                body_run.font.color.rgb = BODY_COLOR
        else:
            body_paragraph = text_frame.paragraphs[0]
            body_paragraph.text = combined_text
            body_paragraph.alignment = PP_ALIGN.LEFT
            body_run = body_paragraph.runs[0]
            body_run.font.name = BODY_FONT_NAME
            body_run.font.size = Pt(SMALL_BODY_FONT_SIZE_POINTS)
            body_run.font.color.rgb = BODY_COLOR

    def render_table(self, slide: Any, table_descriptor: dict[str, Any], x_inches: float, y_inches: float, width_inches: float, max_height_inches: float) -> float:
        """Render an editable PowerPoint table and return the height consumed."""
        row_values = table_descriptor["rows"]
        row_count = len(row_values)
        column_count = max((len(row_value) for row_value in row_values), default=1)
        table_height_inches = min(max_height_inches, max(1.2, 0.36 + 0.26 * row_count))
        if row_count >= 24:
            table_font_size_points = 5
        elif row_count >= 18:
            table_font_size_points = 8
        else:
            table_font_size_points = TABLE_FONT_SIZE_POINTS

        graphic_frame = slide.shapes.add_table(
            row_count,
            column_count,
            Inches(x_inches),
            Inches(y_inches),
            Inches(width_inches),
            Inches(table_height_inches),
        )
        table = graphic_frame.table

        for column_index in range(column_count):
            table.columns[column_index].width = Inches(width_inches / column_count)

        for row_index, row_value in enumerate(row_values):
            for column_index, cell_text in enumerate(row_value):
                cell = table.cell(row_index, column_index)
                cell.text = cell_text
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(245, 248, 251) if row_index == 0 else WHITE_COLOR
                cell.text_frame.word_wrap = True
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER if column_index > 0 else PP_ALIGN.LEFT
                    for run in paragraph.runs:
                        run.font.name = BODY_FONT_NAME
                        run.font.size = Pt(table_font_size_points)
                        run.font.bold = row_index == 0
                        run.font.color.rgb = BODY_COLOR

        if table_descriptor["caption"]:
            caption_y_inches = y_inches + table_height_inches + 0.05
            self.add_textbox(
                slide,
                table_descriptor["caption"],
                x_inches,
                caption_y_inches,
                width_inches,
                0.22,
                font_size_points=CAPTION_FONT_SIZE_POINTS,
                font_color=MUTED_COLOR,
            )
            return table_height_inches + 0.27

        return table_height_inches

    def render_list(self, slide: Any, list_descriptor: dict[str, Any], x_inches: float, y_inches: float, width_inches: float, max_height_inches: float) -> float:
        """Render a bullet or numbered list as editable paragraphs."""
        text_box_shape = slide.shapes.add_textbox(Inches(x_inches), Inches(y_inches), Inches(width_inches), Inches(max_height_inches))
        text_frame = text_box_shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.margin_left = Pt(4)
        text_frame.margin_right = Pt(2)
        text_frame.margin_top = Pt(2)
        text_frame.margin_bottom = Pt(2)

        for item_index, (level, item_text, is_numbered) in enumerate(list_descriptor["items"]):
            paragraph = text_frame.paragraphs[0] if item_index == 0 else text_frame.add_paragraph()
            bullet_prefix = "1. " if is_numbered and level == 0 else "• "
            if level > 0:
                bullet_prefix = "  " * level + "• "
            paragraph.text = f"{bullet_prefix}{item_text}"
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(1)
            for run in paragraph.runs:
                run.font.name = BODY_FONT_NAME
                run.font.size = Pt(BODY_FONT_SIZE_POINTS)
                run.font.color.rgb = BODY_COLOR

        consumed_line_count = sum(max(1, estimate_wrapped_line_count(item_text, 55)) for _, item_text, _ in list_descriptor["items"])
        return min(max_height_inches, 0.15 + consumed_line_count * 0.24)

    def render_paragraph(self, slide: Any, paragraph_text: str, x_inches: float, y_inches: float, width_inches: float, max_height_inches: float, *, font_size_points: int = BODY_FONT_SIZE_POINTS, italic: bool = False, font_color: RGBColor = BODY_COLOR) -> float:
        """Render a plain paragraph and return the height consumed."""
        line_count = estimate_wrapped_line_count(paragraph_text, 70 if width_inches > 10 else 55)
        paragraph_height_inches = min(max_height_inches, max(0.24, 0.12 + line_count * 0.22))
        self.add_textbox(
            slide,
            paragraph_text,
            x_inches,
            y_inches,
            width_inches,
            paragraph_height_inches,
            font_size_points=font_size_points,
            italic=italic,
            font_color=font_color,
        )
        return paragraph_height_inches

    def render_image_grid(self, slide: Any, image_grid_descriptor: dict[str, Any], x_inches: float, y_inches: float, width_inches: float, max_height_inches: float) -> float:
        """Render one or more source figures in an evenly spaced grid."""
        rows = image_grid_descriptor["rows"]
        row_count = len(rows)
        if row_count == 0:
            return 0.0

        gutter_inches = 0.14
        row_heights: list[float] = []
        for row_panels in rows:
            has_labels = any(panel.label_text for panel in row_panels)
            row_heights.append((max_height_inches - gutter_inches * (row_count - 1)) / row_count - (0.2 if has_labels else 0.0))

        current_y_inches = y_inches
        for row_index, row_panels in enumerate(rows):
            panel_width_inches = (width_inches - gutter_inches * (len(row_panels) - 1)) / len(row_panels)
            panel_height_inches = row_heights[row_index]
            for panel_index, image_panel in enumerate(row_panels):
                panel_x_inches = x_inches + panel_index * (panel_width_inches + gutter_inches)
                self.add_contained_picture(slide, image_panel.source_path, panel_x_inches, current_y_inches, panel_width_inches, panel_height_inches)
                if image_panel.label_text:
                    self.render_paragraph(
                        slide,
                        image_panel.label_text,
                        panel_x_inches,
                        current_y_inches + panel_height_inches + 0.02,
                        panel_width_inches,
                        0.26,
                        font_size_points=CAPTION_FONT_SIZE_POINTS,
                        font_color=MUTED_COLOR,
                    )
            current_y_inches += panel_height_inches + gutter_inches
            if any(panel.label_text for panel in row_panels):
                current_y_inches += 0.12

        return current_y_inches - y_inches

    def render_callout_columns(self, slide: Any, columns_descriptor: dict[str, Any], x_inches: float, y_inches: float, width_inches: float, max_height_inches: float) -> float:
        """Render side-by-side callout columns."""
        column_descriptors = columns_descriptor["columns"]
        column_count = len(column_descriptors)
        gutter_inches = 0.18
        column_width_inches = (width_inches - gutter_inches * (column_count - 1)) / column_count
        content_heights = [
            min(max_height_inches, 0.58 + estimate_wrapped_line_count(column_descriptor["title"] + "\n" + column_descriptor["body"], 26) * 0.17)
            for column_descriptor in column_descriptors
        ]
        box_height_inches = min(max_height_inches, max(content_heights, default=1.4))

        for column_index, column_descriptor in enumerate(column_descriptors):
            column_x_inches = x_inches + column_index * (column_width_inches + gutter_inches)
            self.render_callout(slide, column_descriptor, column_x_inches, y_inches, column_width_inches, box_height_inches)

        return box_height_inches

    def render_content_blocks(self, slide: Any, content_blocks: list[dict[str, Any]], frame_number: int, frame_title: str) -> None:
        """Render the normalized block list top-to-bottom with lightweight layout heuristics."""
        if len(content_blocks) == 2 and content_blocks[0]["kind"] == "image_grid" and content_blocks[1]["kind"] == "table":
            image_height_inches = min(4.15, CONTENT_HEIGHT_INCHES * 0.72)
            self.render_image_grid(slide, content_blocks[0], LEFT_MARGIN_INCHES, CONTENT_TOP_INCHES, CONTENT_WIDTH_INCHES, image_height_inches)
            table_y_inches = CONTENT_TOP_INCHES + image_height_inches + 0.08
            table_height_inches = CONTENT_BOTTOM_INCHES - table_y_inches
            self.render_table(slide, content_blocks[1], LEFT_MARGIN_INCHES, table_y_inches, CONTENT_WIDTH_INCHES, table_height_inches)
            return

        if len(content_blocks) == 2 and content_blocks[0]["kind"] == "image_grid" and content_blocks[1]["kind"] == "paragraph":
            image_height_inches = min(4.95, CONTENT_HEIGHT_INCHES * 0.78)
            consumed_height_inches = self.render_image_grid(slide, content_blocks[0], LEFT_MARGIN_INCHES, CONTENT_TOP_INCHES, CONTENT_WIDTH_INCHES, image_height_inches)
            paragraph_y_inches = CONTENT_TOP_INCHES + consumed_height_inches + 0.08
            paragraph_height_inches = CONTENT_BOTTOM_INCHES - paragraph_y_inches
            self.render_paragraph(slide, content_blocks[1]["text"], LEFT_MARGIN_INCHES, paragraph_y_inches, CONTENT_WIDTH_INCHES, paragraph_height_inches)
            return

        if len(content_blocks) == 2 and content_blocks[0]["kind"] == "table" and content_blocks[1]["kind"] == "paragraph":
            paragraph_height_inches = min(0.42, max(0.22, 0.08 + estimate_wrapped_line_count(content_blocks[1]["text"], 110) * 0.10))
            table_height_inches = CONTENT_HEIGHT_INCHES - paragraph_height_inches - 0.22
            self.render_table(slide, content_blocks[0], LEFT_MARGIN_INCHES, CONTENT_TOP_INCHES, CONTENT_WIDTH_INCHES, table_height_inches)
            paragraph_y_inches = CONTENT_TOP_INCHES + table_height_inches + 0.12
            self.render_paragraph(
                slide,
                content_blocks[1]["text"],
                LEFT_MARGIN_INCHES,
                paragraph_y_inches,
                CONTENT_WIDTH_INCHES,
                paragraph_height_inches,
                font_size_points=8,
            )
            return

        current_y_inches = CONTENT_TOP_INCHES
        block_gap_inches = 0.12

        for block_index, content_block in enumerate(content_blocks):
            remaining_height_inches = CONTENT_BOTTOM_INCHES - current_y_inches
            if remaining_height_inches <= 0.18:
                LOGGER.warning(
                    "Slide %s exceeded vertical space and skipped remaining content: %s",
                    frame_number,
                    frame_title,
                )
                break

            content_kind = content_block["kind"]
            consumed_height_inches = 0.0

            if content_kind == "paragraph":
                consumed_height_inches = self.render_paragraph(slide, content_block["text"], LEFT_MARGIN_INCHES, current_y_inches, CONTENT_WIDTH_INCHES, remaining_height_inches)
            elif content_kind == "list":
                consumed_height_inches = self.render_list(slide, content_block, LEFT_MARGIN_INCHES, current_y_inches, CONTENT_WIDTH_INCHES, remaining_height_inches)
            elif content_kind == "callout":
                callout_height_inches = min(remaining_height_inches, max(0.85, 0.48 + estimate_wrapped_line_count(content_block["title"] + "\n" + content_block["body"], 72) * 0.18))
                self.render_callout(slide, content_block, LEFT_MARGIN_INCHES, current_y_inches, CONTENT_WIDTH_INCHES, callout_height_inches)
                consumed_height_inches = callout_height_inches
            elif content_kind == "callout_columns":
                consumed_height_inches = self.render_callout_columns(slide, content_block, LEFT_MARGIN_INCHES, current_y_inches, CONTENT_WIDTH_INCHES, remaining_height_inches)
            elif content_kind == "table":
                consumed_height_inches = self.render_table(slide, content_block, LEFT_MARGIN_INCHES, current_y_inches, CONTENT_WIDTH_INCHES, remaining_height_inches)
            elif content_kind == "image_grid":
                if block_index == len(content_blocks) - 1:
                    allotted_height_inches = remaining_height_inches
                elif any(following_block["kind"] in {"list", "paragraph"} for following_block in content_blocks[block_index + 1:]):
                    allotted_height_inches = min(remaining_height_inches, max(2.2, remaining_height_inches * 0.68))
                else:
                    allotted_height_inches = remaining_height_inches
                consumed_height_inches = self.render_image_grid(slide, content_block, LEFT_MARGIN_INCHES, current_y_inches, CONTENT_WIDTH_INCHES, allotted_height_inches)

            current_y_inches += consumed_height_inches + block_gap_inches

    def build(self) -> None:
        """Generate the editable PowerPoint deck from the Beamer source."""
        pandoc_document = self.load_pandoc_document()
        metadata = pandoc_document["meta"]
        frame_blocks = [block for block in pandoc_document["blocks"] if block["t"] == "Div" and "frame" in block["c"][0][1]]
        total_frame_count = len(frame_blocks)
        LOGGER.info("Building editable PowerPoint with %s slides", total_frame_count)

        presentation = Presentation()
        presentation.slide_width = Inches(SLIDE_WIDTH_INCHES)
        presentation.slide_height = Inches(SLIDE_HEIGHT_INCHES)
        presentation.core_properties.title = inlines_to_text(metadata["title"]["c"])
        presentation.core_properties.author = "\n".join(inlines_to_text(author_item["c"]) for author_item in metadata["author"]["c"])

        self.render_title_slide(presentation, metadata)

        for frame_number, frame_block in enumerate(frame_blocks[1:], start=2):
            frame_title, content_blocks = self.normalize_frame(frame_block)
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            self.add_standard_slide_chrome(slide, frame_title, frame_number, total_frame_count)
            self.render_content_blocks(slide, content_blocks, frame_number, frame_title)

        presentation.save(self.output_presentation_path)
        LOGGER.info("Saved editable PowerPoint to %s", self.output_presentation_path)


def main() -> None:
    """Entry point for building the editable PowerPoint version of the Beamer deck."""
    logging.basicConfig(level=logging.INFO, format="%(asctime)s | %(levelname)s | %(message)s")
    slides_directory = Path(__file__).resolve().parent
    editable_presentation_builder = EditablePresentationBuilder(slides_directory)
    editable_presentation_builder.build()


if __name__ == "__main__":
    main()
