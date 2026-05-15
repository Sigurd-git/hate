import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = pptx.Presentation()
    # Set slide size to 16:9 for modern look
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Colors
    COLOR_DEEP_BLUE = RGBColor(0, 32, 96)
    COLOR_LIGHT_BLUE = RGBColor(0, 112, 192)
    COLOR_GRAY = RGBColor(128, 128, 128)
    COLOR_WHITE = RGBColor(255, 255, 255)

    def add_background_elements(slide):
        # Header bar
        header = slide.shapes.add_shape(
            pptx.enum.shapes.MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = COLOR_DEEP_BLUE
        header.line.width = 0

    def add_title_slide(title_text, subtitle_text):
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
        # Add background
        rect = slide.shapes.add_shape(
            pptx.enum.shapes.MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = COLOR_DEEP_BLUE
        
        # Title
        title = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(2))
        tf = title.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(2))
        tf = subtitle.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle_text
        p.font.size = Pt(24)
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER

    def add_section_slide(section_title):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        rect = slide.shapes.add_shape(
            pptx.enum.shapes.MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = COLOR_LIGHT_BLUE
        
        title = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.33), Inches(1.5))
        p = title.text_frame.paragraphs[0]
        p.text = section_title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER

    def add_content_slide(title_text, points, section_name="", img_path=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_background_elements(slide)
        
        # Section Name (top left)
        if section_name:
            sec_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(4), Inches(0.4))
            p = sec_box.text_frame.paragraphs[0]
            p.text = section_name
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = COLOR_WHITE
            
        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12), Inches(1))
        p = title_box.text_frame.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLOR_DEEP_BLUE
        
        # Text Content (Left side if image exists)
        text_width = Inches(6) if img_path else Inches(11.5)
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), text_width, Inches(5))
        tf = text_box.text_frame
        tf.word_wrap = True
        
        for point in points:
            if isinstance(point, list):
                for sub in point:
                    p = tf.add_paragraph()
                    p.text = "• " + sub
                    p.level = 1
                    p.font.size = Pt(18)
                    p.space_before = Pt(5)
            else:
                p = tf.add_paragraph()
                p.text = "▶ " + point
                p.level = 0
                p.font.size = Pt(22)
                p.font.bold = True
                p.space_before = Pt(15)
        
        # Image (Right side)
        if img_path and os.path.exists(img_path):
            img_left = Inches(7)
            img_top = Inches(2)
            img_height = Inches(4.5)
            # Check aspect ratio
            slide.shapes.add_picture(img_path, img_left, img_top, height=img_height)

    # 1. Title
    add_title_slide(
        "High-level visual representations in the human brain are aligned with large language models",
        "人脑高级视觉表征与大语言模型的对齐\nNature Machine Intelligence | 2025"
    )

    # 2. Section: Background
    add_section_slide("第一部分：背景与核心问题")
    
    add_content_slide("视觉认知的核心挑战", [
        "传统视角：视觉 = 物体识别 (Object Recognition)",
        "现实挑战：自然场景远比物体标签复杂",
        [
            "物体间的空间与语义关系",
            "场景上下文与统计规律",
            "行动可供性 (Affordance)"
        ],
        "科学缺憾：缺乏能够统一量化这些复杂信息的计算框架"
    ], "背景与核心问题", "page_images/page_1.png")

    add_content_slide("计算假设：LLM Embedding 作为桥梁", [
        "为什么是 LLM？",
        [
            "Caption 是对视觉场景的高度概括性语言描述",
            "LLM (如 MPNet) 能学习并编码复杂的上下文关系",
            "语言模型中的统计世界知识可能与视觉系统共享表征几何"
        ],
        "核心路径：图像 → 人工描述 → LLM Embedding ↔ 脑活动"
    ], "背景与核心问题")

    # 3. Section: Exp 1
    add_section_slide("第二部分：实验 1 - 整体映射验证")
    
    add_content_slide("RSA: 表征相似性分析", [
        "研究问题：LLM 空间与脑活动空间的“距离结构”是否一致？",
        "方法流程：",
        [
            "计算图像在 LLM 空间中的 RDM (距离矩阵)",
            "计算图像在脑区内的 RDM (fMRI pattern 差异)",
            "对比两者的相关性 (Searchlight RSA)"
        ],
        "结果：在 Ventral, Lateral, Parietal Streams 发现广泛显著对齐"
    ], "实验 1", "page_images/page_2.png")

    add_content_slide("编码模型 (Encoding Models)", [
        "研究问题：能否用 LLM 特征线性预测单个 Voxel 的响应？",
        "实验细节：",
        [
            "使用 768 维 MPNet 向量作为预测器",
            "Ridge Regression 进行 Voxel-wise 建模",
            "模型表现接近 Noise Ceiling，显示了极强的预测力"
        ]
    ], "实验 1", "page_images/page_2.png")

    # 4. Section: Exp 2
    add_section_slide("第三部分：实验 2 - 语义结构与解码")

    add_content_slide("重现经典选择性 (Selectivity)", [
        "研究问题：映射是否具有功能神经解剖学意义？",
        "预测实验：",
        [
            "输入特定语义句子 (人、地点、食物)",
            "模型预测的脑图重现了 FFA, PPA 等经典脑区分布",
            "证明 Mapping 不仅是统计相关，而是捕捉了功能组织"
        ]
    ], "实验 2", "page_images/page_3.png")

    add_content_slide("语义解码 (Brain Decoding)", [
        "研究问题：从脑活动中能“读出”什么？",
        "方法：",
        [
            "Brain Activity → Linear Decoder → LLM Embedding",
            "Dictionary Search 从 310 万条描述中找回最贴切的 Caption"
        ],
        "结果：精准还原了被试所见的场景内容 (长颈鹿、人群等)"
    ], "实验 2", "page_images/page_3.png")

    # 5. Section: Exp 3
    add_section_slide("第四部分：实验 3 - 机制拆解")

    add_content_slide("排除替代解释：LLM 优在何处？", [
        "对比模型 (Control Models)：",
        [
            "Category-only: 仅保留物体标签信息",
            "Noun-only / Verb-only: 仅提取名词或动词",
            "Word-average: 单词 Embedding 的简单平均"
        ],
        "关键结果：",
        "Full-caption LLM 整合模型在所有 ROI 中表现最优，证明了 Contextual Integration 的核心作用。"
    ], "实验 3", "page_images/page_4.png")

    add_content_slide("词袋 vs. 语篇整合", [
        "结论：高级视觉表征不只是词汇的堆砌 (Bag of words)",
        "证据：单词平均模型显著弱于全文整合模型",
        "启示：脑活动编码的是复杂的“场景叙事”，而非孤立的“物体清单”"
    ], "实验 3")

    # 6. Section: Exp 4
    add_section_slide("第五部分：实验 4 - 计算建模与训练")

    add_content_slide("以 LLM 为目标的视觉模型训练", [
        "计算假设验证：",
        "如果高级视觉目标 ≈ LLM 空间，那么以此为目标的视觉网络应更像脑。",
        "实验设计：",
        [
            "模型：Recurrent CNN (RCNN)",
            "目标：预测图片对应的 LLM Embedding",
            "对比组：同一架构预测分类标签 (Categorical)"
        ]
    ], "实验 4", "page_images/page_5.png")

    add_content_slide("模型评估与跨模型比较", [
        "结果 1：LLM-trained RCNN 表现显著优于类别训练模型",
        "结果 2：与 13 种主流模型 (CLIP, ResNet, etc.) 的全面对齐测试中表现强劲",
        "结论：目标质量 (Semantic Target) 相比数据规模对 Brain-like 更有决定性"
    ], "实验 4", "page_images/page_5.png")

    # 7. Section: Discussion
    add_section_slide("第六部分：综合讨论与展望")

    add_content_slide("核心总结：通向语义表征之路", [
        "1. 视觉系统将图像转换为“类 LLM”的高维语义表征",
        "2. 这种表征涵盖了物体、动作、关系与世界知识",
        "3. 语言模型提供的量化框架为场景理解研究提供了统一工具"
    ], "讨论", "page_images/page_6.png")

    add_content_slide("局限性与未来研究", [
        "任务困扰：识别任务可能引入内部语言叙述 (Captioning)",
        "因果验证：模型相关性与真实神经实现的距离",
        "未来方向：跨物种对比、动态场景整合、具身智能应用"
    ], "讨论")

    # 8. End
    add_title_slide("Thank You!", "Questions & Discussions\n2025.05.04")

    prs.save("visual_llm_expert_v2.pptx")

create_presentation()
