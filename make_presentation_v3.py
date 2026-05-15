import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = pptx.Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    COLOR_DEEP_BLUE = RGBColor(0, 32, 96)
    COLOR_LIGHT_BLUE = RGBColor(0, 112, 192)
    COLOR_WHITE = RGBColor(255, 255, 255)

    def add_background_elements(slide):
        header = slide.shapes.add_shape(
            pptx.enum.shapes.MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = COLOR_DEEP_BLUE
        header.line.width = 0

    def add_title_slide(title_text, subtitle_text):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        rect = slide.shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        rect.fill.solid(); rect.fill.fore_color.rgb = COLOR_DEEP_BLUE
        title = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(2))
        p = title.text_frame.paragraphs[0]
        p.text = title_text; p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = COLOR_WHITE; p.alignment = PP_ALIGN.CENTER
        subtitle = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(2))
        p = subtitle.text_frame.paragraphs[0]
        p.text = subtitle_text; p.font.size = Pt(24); p.font.color.rgb = COLOR_WHITE; p.alignment = PP_ALIGN.CENTER

    def add_section_slide(section_title):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        rect = slide.shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        rect.fill.solid(); rect.fill.fore_color.rgb = COLOR_LIGHT_BLUE
        title = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.33), Inches(1.5))
        p = title.text_frame.paragraphs[0]
        p.text = section_title; p.font.size = Pt(48); p.font.bold = True; p.font.color.rgb = COLOR_WHITE; p.alignment = PP_ALIGN.CENTER

    def add_content_slide(title_text, points, section_name="", img_path=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_background_elements(slide)
        if section_name:
            sec_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(4), Inches(0.4))
            p = sec_box.text_frame.paragraphs[0]; p.text = section_name; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = COLOR_WHITE
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12), Inches(1))
        p = title_box.text_frame.paragraphs[0]; p.text = title_text; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = COLOR_DEEP_BLUE
        
        text_width = Inches(6.5) if img_path else Inches(12)
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), text_width, Inches(5))
        tf = text_box.text_frame; tf.word_wrap = True
        for point in points:
            if isinstance(point, list):
                for sub in point:
                    p = tf.add_paragraph(); p.text = "• " + sub; p.level = 1; p.font.size = Pt(16); p.space_before = Pt(3)
            else:
                p = tf.add_paragraph(); p.text = "▶ " + point; p.level = 0; p.font.size = Pt(20); p.font.bold = True; p.space_before = Pt(12)
        
        if img_path and os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(7.2), Inches(1.5), height=Inches(5.5))

    # --- PPT CONSTRUCTION (Total 30 Slides) ---
    
    # 1. Title
    add_title_slide("High-level visual representations in the human brain are aligned with large language models", "自然场景高级视觉表征与大语言模型的对齐\n汇报人：[您的姓名]\n日期：2026年5月4日")
    
    # 2-6. Background & Core Question
    add_section_slide("第一部分：研究背景与核心假设")
    add_content_slide("视觉认知的“物体识别”范式", ["传统研究主要关注视觉系统如何从像素中识别物体分类","关键成果：发现 FFA, PPA, EBA 等类别选择性区域","局限：自然场景理解不等于识别孤立的物体标签"], "研究背景")
    add_content_slide("自然场景的复杂语义挑战", ["场景包含丰富的统计规律与物体间关系","例如：‘狗在船上’（空间关系） vs. ‘狗’ + ‘船’（物体列表）","缺乏能够统一建模这些复杂高层信息的量化框架"], "研究背景", "page_images/page_1.png")
    add_content_slide("大语言模型 (LLM) 的语义潜能", ["LLM (如 MPNet) 能将场景描述 (Caption) 编码为高维连续空间","其优势在于：捕捉词与词之间的上下文互动与世界知识","假设：这种 LLM 语义空间可能与脑的高级视觉空间对齐"], "研究背景")
    add_content_slide("论证逻辑：四大实验设计", ["Exp 1: 验证 LLM Embedding 与脑活动的整体相关性","Exp 2: 验证该映射的功能特异性与可解码性","Exp 3: 拆解 LLM 的优势来源：排除替代解释","Exp 4: 训练以 LLM 为目标的视觉模型进行闭环验证"], "研究背景")

    # 7-11. Exp 1
    add_section_slide("第二部分：实验 1 - 整体映射验证")
    add_content_slide("NSD 数据集与方法论", ["Natural Scenes Dataset (NSD): 8被试, 7T fMRI","输入：每张图像配有 5 条人工 Caption","特征：将 Caption 输入 LLM 提取 Embedding"], "实验 1", "page_images/page_1.png")
    add_content_slide("RSA (表征相似性分析) 设计", ["计算图像在 LLM 空间中的距离结构 (Model RDM)","计算图像在脑区 pattern 中的差异结构 (Brain RDM)","对比两者的几何相似性"], "实验 1")
    add_content_slide("结果：广泛的层级对齐", ["在腹侧、外侧及顶侧流的高级视觉区域发现显著对齐","证明了 LLM 语义空间与脑视觉表征的强关联"], "实验 1", "page_images/page_2.png")
    add_content_slide("Voxel-wise Encoding 验证", ["使用 LLM 特征线性预测每个 Voxel 的 beta 响应","预测准确度接近被试间一致性上限 (Noise Ceiling)","显示了极强的层级语义刻画能力"], "实验 1", "page_images/page_2.png")

    # 12-16. Exp 2
    add_section_slide("第三部分：实验 2 - 语义结构与功能映射")
    add_content_slide("功能选择性重现 (Selectivity Contrast)", ["利用编码模型预测特定语义句子诱发的脑活动","对比：‘人’ vs. ‘场景’；‘食物’ vs. ‘人’","准确重现了经典的功能区分布 (FFA, PPA, Food Area)"], "实验 2", "page_images/page_3.png")
    add_content_slide("语义解码：从脑活动还原场景", ["设计线性解码器：Brain Pattern → Predicted Embedding","从 3.1M 库中搜索最近邻 Caption","成功还原了被试所见的复杂图像内容"], "实验 2", "page_images/page_3.png")
    add_content_slide("解码结果示例与分析", ["示例：长颈鹿、桌边人群、停车位等场景还原","结论：被预测出的不只是低级特征，而是可语言化的场景语义"], "实验 2")
    add_content_slide("对齐的结构化意义", ["证明了 Mapping 具有明确的功能神经解剖学基础","验证了视觉表征空间具有可解码的、稳定的语义结构"], "实验 2")

    # 17-21. Exp 3
    add_section_slide("第四部分：实验 3 - 机制拆解与排除解释")
    add_content_slide("对照模型 A：排除“仅物体类别”", ["问题：对齐是否仅因为 LLM 包含了物体标签？","结果：Full Caption 显著优于 Category-only 模型","结论：关系与动作信息对脑对齐有关键贡献"], "实验 3", "page_images/page_4.png")
    add_content_slide("对照模型 B：排除“仅名词/动词列表”", ["问题：是否只是因为 LLM 认识更多词汇？","结果：整合的 Full Caption 优于名词或动词的简单集合","进一步支持：脑表征是整合的场景语义空间"], "实验 3", "page_images/page_4.png")
    add_content_slide("对照模型 C：排除“词袋平均” (关键环节)", ["问题：词序和上下文整合是否必要？","结果：Full-sentence Embedding 优于单词平均模型","启示：脑不仅编码‘有什么’，还编码‘如何互动’"], "实验 3", "page_images/page_4.png")
    add_content_slide("机制总结：上下文整合的核心作用", ["LLM 的优势源自其对复杂语义关系的上下文建模能力","这也正是高级视觉皮层区别于早期视觉区的核心特质"], "实验 3")

    # 22-26. Exp 4
    add_section_slide("第五部分：实验 4 - 视觉模型的计算建模")
    add_content_slide("以 LLM 为目标的视觉模型训练", ["计算假设：如果 LLM 空间是脑的目标，则以此为训练目标的模型更像脑","架构：Recurrent CNN (RCNN)","目标：图片 → 预测 LLM Caption Embedding"], "实验 4", "page_images/page_5.png")
    add_content_slide("LLM 目标 vs. 分类标签目标", ["同一架构下：LLM-target 模型比 Category-target 模型更 Brain-like","证明了语义目标相比单纯分类标签的优越性"], "实验 4", "page_images/page_5.png")
    add_content_slide("跨模型性能大比拼 (SOTA 比较)", ["与 CLIP, ResNet 等 13 种主流模型在 NSD 上的对齐度对比","LLM-trained RCNN 在数据量极少的情况下达到顶尖性能"], "实验 4", "page_images/page_5.png")
    add_content_slide("类别信息的读出实验", ["从 LLM-trained 模型中可成功读出物体分类信息","说明 LLM 表征包含了类别，但提供了更丰富的特征结构"], "实验 4")

    # 27-30. Discussion
    add_section_slide("第六部分：综合讨论与未来展望")
    add_content_slide("视觉认知的范式迁移", ["从‘物体清单’到‘语义场景叙事’","LLM Embedding 为这种统一建模提供了强大的数学工具","视觉与语言在高级语义空间上达成对齐"], "讨论", "page_images/page_6.png")
    add_content_slide("局限性：内部言语化的困扰", ["任务设置可能诱发被试的内部描述 (Captioning)","未来需在无意识或快速呈现任务下进一步验证对齐性"], "讨论")
    add_content_slide("结论与展望", ["视觉系统不仅在看，更是在‘叙述’世界的语义","这一对齐发现为通向通用智能的视觉语义化提供了脑科学依据"], "讨论")
    add_title_slide("Thank You!", "Q&A | 谢 谢 聆 听\nNature Machine Intelligence 2025")

    prs.save("visual_llm_expert_full_30.pptx")

create_presentation()
print("Generated 30 slides in visual_llm_expert_full_30.pptx")
